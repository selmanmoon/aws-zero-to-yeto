"""
PPTX to PDF Converter Lambda Function
=====================================
Triggered by S3 .pptx uploads in the 'pptxs/' prefix.
Uses python-pptx to extract content and reportlab to generate PDF summaries.
"""

import json
import os
import boto3
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# Initialize S3 client
s3_client = boto3.client('s3')

def extract_text_from_shape(shape):
    """Extract text from a shape if it has a text frame."""
    text_parts = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = ""
            for run in paragraph.runs:
                paragraph_text += run.text
            if paragraph_text.strip():
                text_parts.append(paragraph_text.strip())
    return text_parts

def extract_slide_content(slide, slide_number):
    """Extract all text content from a slide."""
    slide_content = {
        'slide_number': slide_number,
        'title': None,
        'content': []
    }
    
    # Try to get the slide title
    if slide.shapes.title:
        slide_content['title'] = slide.shapes.title.text.strip() if slide.shapes.title.text else f"Slide {slide_number}"
    else:
        slide_content['title'] = f"Slide {slide_number}"
    
    # Extract text from all shapes
    for shape in slide.shapes:
        # Skip the title shape as we already captured it
        if slide.shapes.title and shape == slide.shapes.title:
            continue
        
        text_parts = extract_text_from_shape(shape)
        slide_content['content'].extend(text_parts)
    
    return slide_content

def create_pdf_from_content(slides_content, original_filename):
    """Create a PDF document from extracted slide content using ReportLab."""
    buffer = BytesIO()
    
    # Create the PDF document
    doc = SimpleDocTemplate(
        buffer,
        pagesize=letter,
        rightMargin=inch,
        leftMargin=inch,
        topMargin=inch,
        bottomMargin=inch
    )
    
    # Define styles
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=TA_CENTER,
        textColor='#1a365d'
    )
    
    slide_title_style = ParagraphStyle(
        'SlideTitle',
        parent=styles['Heading2'],
        fontSize=16,
        spaceBefore=20,
        spaceAfter=12,
        textColor='#2d3748',
        borderWidth=0,
        borderPadding=0,
        borderColor='#e2e8f0',
        backColor='#f7fafc'
    )
    
    content_style = ParagraphStyle(
        'ContentStyle',
        parent=styles['Normal'],
        fontSize=11,
        spaceBefore=6,
        spaceAfter=6,
        leftIndent=20,
        textColor='#4a5568'
    )
    
    # Build the document content
    story = []
    
    # Document title
    doc_title = original_filename.replace('.pptx', '').replace('_', ' ').title()
    story.append(Paragraph(f"📄 {doc_title}", title_style))
    story.append(Paragraph(f"Generated on: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')} UTC", styles['Normal']))
    story.append(Paragraph(f"Total Slides: {len(slides_content)}", styles['Normal']))
    story.append(Spacer(1, 30))
    
    # Add each slide's content
    for slide_data in slides_content:
        # Slide title
        slide_title = f"Slide {slide_data['slide_number']}: {slide_data['title']}"
        story.append(Paragraph(slide_title, slide_title_style))
        
        # Slide content
        if slide_data['content']:
            for text in slide_data['content']:
                # Escape special characters for ReportLab
                safe_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                story.append(Paragraph(f"• {safe_text}", content_style))
        else:
            story.append(Paragraph("<i>No text content on this slide</i>", content_style))
        
        story.append(Spacer(1, 15))
    
    # Build PDF
    doc.build(story)
    
    buffer.seek(0)
    return buffer

def lambda_handler(event, context):
    """
    Main Lambda handler - processes PPTX files and generates PDF summaries.
    
    Triggered by S3 ObjectCreated events in the 'pptxs/' prefix.
    """
    print("🚀 Converter Lambda triggered")
    print(f"Event: {json.dumps(event, indent=2)}")
    
    try:
        # Get bucket and key from S3 event
        record = event['Records'][0]
        bucket_name = record['s3']['bucket']['name']
        object_key = record['s3']['object']['key']
        
        print(f"📁 Processing: s3://{bucket_name}/{object_key}")
        
        # Validate file extension
        if not object_key.lower().endswith('.pptx'):
            print(f"⚠️ Skipping non-PPTX file: {object_key}")
            return {
                'statusCode': 200,
                'body': json.dumps('Skipped: Not a PPTX file')
            }
        
        # Download PPTX from S3
        print("📥 Downloading PPTX from S3...")
        response = s3_client.get_object(Bucket=bucket_name, Key=object_key)
        pptx_content = response['Body'].read()
        
        # Parse PPTX
        print("📊 Parsing PPTX content...")
        pptx_buffer = BytesIO(pptx_content)
        presentation = Presentation(pptx_buffer)
        
        # Extract content from all slides
        slides_content = []
        for idx, slide in enumerate(presentation.slides, start=1):
            slide_data = extract_slide_content(slide, idx)
            slides_content.append(slide_data)
            print(f"  ✅ Processed slide {idx}: {slide_data['title']}")
        
        slide_count = len(slides_content)
        print(f"📄 Total slides extracted: {slide_count}")
        
        # Generate PDF
        print("🔄 Generating PDF...")
        original_filename = os.path.basename(object_key)
        pdf_buffer = create_pdf_from_content(slides_content, original_filename)
        
        # Prepare output paths
        base_name = original_filename.replace('.pptx', '').replace('.PPTX', '')
        pdf_key = f"pdfs/{base_name}.pdf"
        metadata_key = f"metadata/{base_name}.json"
        
        # Upload PDF to S3
        print(f"📤 Uploading PDF to: {pdf_key}")
        s3_client.put_object(
            Bucket=bucket_name,
            Key=pdf_key,
            Body=pdf_buffer.getvalue(),
            ContentType='application/pdf'
        )
        
        # Create and upload metadata JSON
        metadata = {
            'pdf_name': f"{base_name}.pdf",
            'original_name': original_filename,
            'slide_count': slide_count,
            'create_date': datetime.utcnow().isoformat() + 'Z',
            'pdf_path': pdf_key,
            'source_path': object_key
        }
        
        print(f"📤 Uploading metadata to: {metadata_key}")
        s3_client.put_object(
            Bucket=bucket_name,
            Key=metadata_key,
            Body=json.dumps(metadata, indent=2),
            ContentType='application/json'
        )
        
        print("✅ Conversion completed successfully!")
        
        return {
            'statusCode': 200,
            'body': json.dumps({
                'message': 'PPTX converted successfully',
                'pdf_path': pdf_key,
                'metadata_path': metadata_key,
                'slide_count': slide_count
            })
        }
        
    except Exception as e:
        print(f"❌ Error processing PPTX: {str(e)}")
        import traceback
        traceback.print_exc()
        
        return {
            'statusCode': 500,
            'body': json.dumps({
                'error': str(e),
                'message': 'Failed to convert PPTX'
            })
        }
